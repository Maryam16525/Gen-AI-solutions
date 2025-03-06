from openai import AzureOpenAI
from dotenv import load_dotenv
import os
from azure.identity import DefaultAzureCredential
from azure.storage.blob import BlobServiceClient
from azure.kusto.data import KustoClient, KustoConnectionStringBuilder
from tenacity import retry, wait_random_exponential, stop_after_attempt
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader
from azure.kusto.ingest import KustoStreamingIngestClient, IngestionProperties
from azure.kusto.data.helpers import dataframe_from_result_table

import streamlit as st
import zipfile
import time

from azure.kusto.data.exceptions import KustoServiceError

import pandas as pd
from pptx import Presentation  # For PowerPoint files
from docx import Document  # For Word documents
from openpyxl import load_workbook
import json

from pptx import Presentation  # For PowerPoint files
from docx import Document  # For Word documents
from openpyxl import load_workbook
import os

# Load environment variables from .env file
load_dotenv()

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
OPENAI_DEPLOYMENT_ENDPOINT = os.getenv("OPENAI_DEPLOYMENT_ENDPOINT")
api_version = api_version="2023-12-01-preview"
OPENAI_ADA_EMBEDDING_DEPLOYMENT_NAME = "text-embedding-ada-002"
OPENAI_GPT4_DEPLOYMENT_NAME = "gpt-4o"
KUSTO_DATABASE = "GenAI_eventhouse"
KUSTO_TABLE = "bookEmbeddings"

container_name = "container1"
blob_name = "temp_batch_0.csv"

KUSTO_URI=os.getenv("KUSTO_URI")

connection_string =os.getenv("connection_string")


# Initialize OpenAI client
client = AzureOpenAI(
    azure_endpoint=OPENAI_DEPLOYMENT_ENDPOINT,
    api_key=OPENAI_API_KEY,
    api_version="2023-12-01-preview"
)
# Create Default Azure Credential
credential = DefaultAzureCredential()
token = credential.get_token("https://kusto.kusto.windows.net/.default").token
# Create BlobServiceClient
connection_string=os.getenv("AZURE_BLOB_CONNECTION_STRING")
blob_service_client = BlobServiceClient.from_connection_string(connection_string)

print(f"Connection string: {connection_string}")  # This should not be None
# Upload CSV to blob storage
def upload_to_blob(file_path, blob_name):
    blob_client = blob_service_client.get_blob_client(container=container_name, blob=blob_name)
    with open(file_path, "rb") as data:
        blob_client.upload_blob(data, overwrite=True)
    print(f"File {file_path} uploaded to Blob Storage.")
# Kusto connection setup
kcsb = KustoConnectionStringBuilder.with_aad_user_token_authentication(KUSTO_URI, token)
client_KUSTO = KustoClient(kcsb)
# Function to generate embeddings using OpenAI
@retry(wait=wait_random_exponential(min=1, max=20), stop=stop_after_attempt(6))
# Function to call OpenAI for chat completions
def call_openAI(text):
    response = client.chat.completions.create(
        model=OPENAI_GPT4_DEPLOYMENT_NAME,
        messages = text,
        temperature=0
    )
    return response.choices[0].message.content
# Load and split PDFs
splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=30)


def load_documents_in_batches(directory, chatbot_name, splitter, batch_size=10, placeholder=None):
    df = pd.DataFrame(columns=['document_name', 'content', 'embedding', 'PartitionKey'])
    batch_files = []
    # Get the total number of files for progress tracking
    total_files = sum(len(files) for _, _, files in os.walk(directory))
    processed_files = 0
    for root, dirs, files in os.walk(directory):
        for file_name in files:
            file_path = os.path.join(root, file_name)
            batch_files.append(file_path)
            # Process the batch when it reaches the specified batch size
            if len(batch_files) >= batch_size:
                process_batch(batch_files, df, chatbot_name, splitter)
                processed_files += len(batch_files)
                if placeholder:
                    placeholder.text(f"Processed {processed_files} of {total_files} files...")
                batch_files = []  # Clear the batch after processing
    # Process any remaining files in the last batch
    if batch_files:
        process_batch(batch_files, df, chatbot_name, splitter)
        processed_files += len(batch_files)
        if placeholder:
            placeholder.text(f"Processed {processed_files} of {total_files} files...")
    return df
def process_batch(batch_files, df, chatbot_name, splitter):
    # Initialize the progress bar
    progress_bar = st.progress(0)
    
    total_files = len(batch_files)
    
    for idx, file_path in enumerate(batch_files):
        try:
            if file_path.endswith('.pdf'):
                st.info(f"Processing PDF file: {file_path}")  # Update UI status
                loader = PyPDFLoader(file_path)
                pages = loader.load_and_split(text_splitter=splitter)
                for page in pages:
                    df.loc[len(df.index)] = [os.path.basename(file_path), page.page_content, "", chatbot_name]
                st.success(f"Successfully processed PDF file: {file_path}")  # Success message
            
            elif file_path.endswith('.csv'):
                st.info(f"Processing CSV file: {file_path}")  # Update UI status
                csv_data = pd.read_csv(file_path)
                content = csv_data.to_string(index=False)
                df.loc[len(df.index)] = [os.path.basename(file_path), content, "", chatbot_name]
                st.success(f"Successfully processed CSV file: {file_path}")  # Success message
            
            elif file_path.endswith(('.ppt', '.pptx')):
                st.info(f"Processing PPT file: {file_path}")  # Update UI status
                ppt = Presentation(file_path)
                content = ""
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            content += shape.text + "\n"
                df.loc[len(df.index)] = [os.path.basename(file_path), content, "", chatbot_name]
                st.success(f"Successfully processed PPT file: {file_path}")  # Success message
            
            elif file_path.endswith(('.doc', '.docx')):
                st.info(f"Processing Word file: {file_path}")  # Update UI status
                doc = Document(file_path)
                content = "\n".join([para.text for para in doc.paragraphs])
                df.loc[len(df.index)] = [os.path.basename(file_path), content, "", chatbot_name]
                st.success(f"Successfully processed Word file: {file_path}")  # Success message
            
            elif file_path.endswith('.xlsx'):
                st.info(f"Processing Excel file: {file_path}")  # Update UI status
                workbook = load_workbook(file_path)
                content = ""
                for sheet in workbook.sheetnames:
                    ws = workbook[sheet]
                    for row in ws.iter_rows(values_only=True):
                        content += "\t".join([str(cell) for cell in row]) + "\n"
                df.loc[len(df.index)] = [os.path.basename(file_path), content, "", chatbot_name]
                st.success(f"Successfully processed Excel file: {file_path}")  # Success message
            
            # Update the progress bar
            progress = (idx + 1) / total_files
            progress_bar.progress(progress)
        
        except Exception as e:
            st.error(f"Error loading {file_path}: {e}")  # Error message

    # Complete the progress bar
    progress_bar.progress(1.0)
# Step 1: Alter the table to add the PartitionKey column
def generate_embeddings(text, document_name, verbose=False):
    txt = text.replace("\n", " ")
    try:
        if verbose:
            st.info(f"Generating embeddings for document: {document_name}...")
        embedding = client.embeddings.create(input=[txt], model=OPENAI_ADA_EMBEDDING_DEPLOYMENT_NAME).data[0].embedding
        if verbose:
            st.success(f"Embeddings generated for document: {document_name}.")
        return embedding
    except Exception as e:
        st.error(f"Error generating embedding for {document_name}: {e}")
        return None
# Ingest PDFs into Kusto and upload via Blob Storage
def ingest_data_to_kusto(df, chatbot_name, chunk_size=10, placeholder=None, file_names=None):
    # Add PartitionKey and generate embeddings
    df["PartitionKey"] = chatbot_name
    df["embedding"] = df.apply(lambda row: generate_embeddings(row['content'], row['document_name']), axis=1)
    
    # Notify user that embedding generation is complete
    if placeholder:
        placeholder.text("Embeddings generated successfully.")
    # Filter out rows where embedding generation failed
    df = df[df['embedding'].notnull()]
    ingest_client = KustoStreamingIngestClient(kcsb)
    ingestion_props = IngestionProperties(
        database=KUSTO_DATABASE,
        table=KUSTO_TABLE,
        data_format='csv'
    )
    
    # Notify user that upload to Blob Storage is starting
    if placeholder:
        placeholder.text("Preparing to upload data to Blob Storage...")
    # Upload the entire DataFrame to Blob Storage if needed
    df.to_csv("temp.csv", index=False)
    upload_to_blob("temp.csv", blob_name)  # Upload the CSV to Blob Storage
    
    # Notify user about the file upload
    if placeholder and file_names:
        placeholder.text(f"Uploaded file to Blob Storage: {blob_name}")
    # Notify user that ingestion is starting
    if placeholder:
        placeholder.text("Starting ingestion into Kusto...")
        placeholder.text("Preparing to ingest data. This may take a moment...")
    # Creative notifications during ingestion preparation
    if placeholder:
        placeholder.text("Gathering the data wizards... 🧙‍♂️")
        placeholder.text("Aligning the stars for optimal ingestion... 🌌")
        placeholder.text("Charging the ingestion crystals... 🔮")
        placeholder.text("Almost there... Just a few more moments! ⏳")
    # Ingest in chunks
    total_batches = (len(df) + chunk_size - 1) // chunk_size  # Calculate total batches
    progress_bar = st.progress(0)  # Initialize progress bar
    # Display the files being ingested
    if file_names:
        placeholder.text(f"Ingesting files: {', '.join(file_names)}")
    # Notify user about the chunking process
    if placeholder:
        placeholder.text(f"Total batches to ingest: {total_batches}. Starting the ingestion process...")
    for start in range(0, len(df), chunk_size):
        chunk = df.iloc[start:start + chunk_size]
        chunk.to_csv("temp_chunk.csv", index=False)
        
        with open("temp_chunk.csv", "r") as file:
            ingest_client.ingest_from_dataframe(chunk, ingestion_properties=ingestion_props)
        
        # Update the placeholder with the current batch status
        if placeholder:
            placeholder.text(f"🚀 Ingesting batch {start // chunk_size + 1} of {total_batches}...")
            placeholder.text("Hold tight! Your data is on its way! ✈️")
        
        # Update the progress bar
        progress = (start // chunk_size + 1) / total_batches
        progress_bar.progress(progress)
    if placeholder:
        placeholder.success("🎉 Ingestion complete! Your chatbot is ready to assist you!")
# Hybrid Query - Keyword and Cosine Similarity Search in Kusto
def get_answer_from_eventhouse(question, chatbot_name, nr_of_answers=1, search_term=None):
    # Generate embedding for the question
    searchedEmbedding = generate_embeddings(question, "search_query", verbose=False)
    # Create the Kusto query with both exact search and cosine similarity
    kusto_query = f"""
        let chatbot_name = '{chatbot_name}';
        let search_term = '{search_term if search_term else ""}';  
        let searchedEmbedding = dynamic({str(searchedEmbedding)});
        let nr_of_answers = {nr_of_answers};
        
        {KUSTO_TABLE}
        | where PartitionKey == chatbot_name
        | where content contains search_term  // Keyword/phrase search if a term is provided
        | extend similarity = series_cosine_similarity(searchedEmbedding, embedding)
        | top nr_of_answers by similarity desc
        | where similarity > 0.5  // Adjust threshold if necessary
    """
    
    # Print the Kusto query for debugging
    print("Kusto Query Executed:")
    print(kusto_query)
    response = client_KUSTO.execute(KUSTO_DATABASE, kusto_query)
    kustoDf = dataframe_from_result_table(response.primary_results[0])
    return kustoDf  # Ensure you return the DataFrame for further processing
# Page configuration for Chatbot Builder
# Page configuration
page_config = {
    "name": "Chatbot Builder",                 # name to be used in menu and in the browser window
    "icon": ":material/explore:",              # icon to be used in the menu
    "public": True,                            # should page be on the production site. Default is False => dev only
    "layout": "centered",                          # page layout. ["centered" or "wide"]. Default is centered
    "title": "Self-Service Chatbot Builder",   # title to be shown on page with st.title(). Default is icon + name. If set to None, nothing is shown.
  # list of roles that can access the page. If empty, everyone can access.
}
# import json
def list_json_blobs():
    container_client = blob_service_client.get_container_client(container_name)
    json_files = []
    blob_list = container_client.list_blobs()
    for blob in blob_list:
        if blob.name.endswith('.json'):
            json_files.append(blob.name)
    return json_files
def read_json_blob(blob_name):
    blob_client = blob_service_client.get_blob_client(container=container_name, blob=blob_name)
    blob_data = blob_client.download_blob()
    return json.loads(blob_data.readall())
# Page configuration
page_config = {
    "name": "Chatbot Builder",  # name to be used in menu and in the browser window
    "public": False,  # should page be on the production site. Default is dev only
}
# Add Font Awesome for icons
st.markdown('<link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.0.0-beta3/css/all.min.css">', unsafe_allow_html=True)
# Disclaimer
st.markdown("""
    <div style='background-color: #fff3cd; padding: 10px; border-radius: 5px; border: 1px solid #ffeeba; text-align: center;'>
        <strong style="color: #856404;">Disclaimer:</strong> 
        <span style="color: #856404;">The chatbots you create are accessible to the entire organization. While the names of the chatbots are visible and users can interact with them, they will not have access to the configurations, including instructions and data. 
        Please contact us if you need to secure your Chatbot.</span>
    </div>
""", unsafe_allow_html=True)
# Centered container with fixed width
st.markdown("""
    <div style='max-width: 400px; margin: auto;'>  <!-- Reduced max-width -->
""", unsafe_allow_html=True)

# Chatbot Configuration and Interaction Section
st.markdown('<h2 style="color: #4CAF50;">Chatbot Interaction <i class="fas fa-comments" style="color: #4CAF50;"></i></h2>', unsafe_allow_html=True)
# Load existing chatbots from blob storage
chatbot_files = list_json_blobs()
if chatbot_files:
    selected_chatbot = st.selectbox("Select a Chatbot", chatbot_files, key="select_chatbot_unique")
else:
    st.warning("No chatbots available. Please create one.")
# Create two columns for user input with a gap
cols = st.columns([1, 1])  # Equal width columns
with cols[0]:
    user_input = st.text_input("Ask your question:", key="user_input", max_chars=100, 
                                placeholder="Type your question...")
with cols[1]:
    search_term = st.text_input("Search Term (optional):", key="search_term_input", max_chars=50, 
                                 placeholder="Type a search term...")
# If "Ask Chatbot" button is clicked
if st.button("Ask Chatbot", key="ask_chatbot_button"):
    if selected_chatbot:
        config = read_json_blob(selected_chatbot)
        
        if not user_input:
            st.warning("Please enter a question.")
        else:
            nr_of_answers = 2
            with st.spinner("Searching for answers..."):
                answers_df = get_answer_from_eventhouse(user_input, config['name'], nr_of_answers, search_term)
                
                if answers_df.empty:
                    st.warning("No answers found.")
                else:
                    document_names = [row['document_name'] for index, row in answers_df.iterrows()]
                    source_document = document_names[0]
                    answer = " ".join(answers_df['content'].tolist())
                    st.subheader("Source Document:")
                    st.text(source_document)
                    prompt = f'Question: {user_input}\nInformation: {answer}'
                    
                    messages = [
                        {"role": "system", "content": config['instructions']},
                        {"role": "user", "content": prompt}
                    ]
                    
                    result = call_openAI(messages)
                    st.write("Chatbot's response:", result)
    else:
        st.warning("Please select a chatbot.")
# Create New Chatbot Section
st.markdown('<h2 style="color: #2196F3;">Create New Chatbot <i class="fas fa-cogs" style="color: #2196F3;"></i></h2>', unsafe_allow_html=True)
# Create two columns for new chatbot input with a gap
cols = st.columns([1, 1])  # Equal width columns
with cols[0]:
    new_chatbot_name = st.text_input("New Chatbot Name:", "")
with cols[1]:
    new_chatbot_personality = st.text_input("Chatbot Personality:", "Friendly and helpful.")
new_chatbot_instructions = st.text_area("Chatbot Instructions:", "You are a helpful assistant.")

# Upload ZIP file and build chatbot
uploaded_zip = st.file_uploader("Upload a ZIP file containing documents", type=["zip"], key="zip_uploader")

# Create and Build Chatbot Button
if st.button("Build Chatbot"):
    if new_chatbot_name:
        # Create new chatbot configuration
        new_chatbot_config = {
            "name": new_chatbot_name,
            "personality": new_chatbot_personality,
            "instructions": new_chatbot_instructions
        }
        blob_client = blob_service_client.get_blob_client(container=container_name, blob=f"{new_chatbot_name}.json")
        blob_client.upload_blob(json.dumps(new_chatbot_config), overwrite=True)
        st.success(f"Chatbot '{new_chatbot_name}' created successfully!")

        # Check if a ZIP file is uploaded
        if uploaded_zip:
            with zipfile.ZipFile(uploaded_zip, 'r') as zip_ref:
                zip_ref.extractall("temp_folder")
                extracted_files = zip_ref.namelist()
                st.success("Files extracted from the ZIP!")
                st.write("Files in the ZIP:")
                for file_name in extracted_files:
                    st.write(file_name)

            # Building the chatbot
            try:
                placeholder = st.empty()
                placeholder.text("Building chatbot, please wait...")
                total_files = len(extracted_files)
                placeholder.text("Loading documents...")
                for i, file_name in enumerate(extracted_files):
                    placeholder.text(f"Processing file {i + 1} of {total_files}: {file_name}...")
                    time.sleep(1)  # Simulating processing time
                
                placeholder.text("All files have been loaded successfully.")
                placeholder.text("Generating embeddings...")
                df = load_documents_in_batches("temp_folder", new_chatbot_name, splitter)
                
                if df is not None and not df.empty:
                    placeholder.text("Embeddings generated successfully.")
                    placeholder.text("Preparing data for ingestion...")
                    placeholder.text("Starting ingestion into Kusto...")
                    try:
                        ingest_data_to_kusto(df, new_chatbot_name, chunk_size=10, placeholder=placeholder, file_names=extracted_files)  
                        placeholder.success(f"Chatbot '{new_chatbot_name}' built and data ingested into Kusto successfully!")
                    except Exception as e:
                        placeholder.error(f"Error during ingestion: {e}")
                else:
                    placeholder.warning("No documents found to ingest.")
            except Exception as e:
                placeholder.error(f"Error building chatbot: {e}")
        else:
            st.warning("Please upload a ZIP file to build the chatbot.")
    else:
        st.warning("Please enter a name for the new chatbot.")

# Close the centered container
st.markdown("</div>", unsafe_allow_html=True)
